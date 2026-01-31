"""
教学资源文件解析工具
支持从OSS下载并解析PDF、Word、PPT等文件格式
"""
import os
import tempfile
import logging
from typing import Optional
from app.models.teaching_resource import TeachingResource
from app.utils.oss_client import oss_client

logger = logging.getLogger(__name__)

# 文本长度限制
MAX_TEXT_LENGTH = 20000


async def download_and_parse_resource(resource: TeachingResource) -> str:
    """
    从OSS下载资源并解析文本内容
    
    Args:
        resource: 教学资源对象
        
    Returns:
        str: 解析后的文本内容
    """
    try:
        # 1. 创建临时文件
        suffix = get_file_suffix(resource.resource_type, resource.original_filename)
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        temp_path = temp_file.name
        temp_file.close()
        
        # 2. 从OSS下载文件
        logger.info(f"开始下载资源: {resource.id}, 类型: {resource.resource_type}")
        
        # 优先使用PDF版本（如果已转换）
        if resource.pdf_path and resource.pdf_conversion_status == 'success':
            download_path = resource.pdf_path
            resource_type = 'pdf'
        else:
            download_path = resource.file_path
            resource_type = resource.resource_type
        
        # 下载文件
        if not oss_client.enabled:
            raise Exception("OSS未配置，无法下载文件")
        
        try:
            # 使用oss2的get_object方法下载文件
            result = oss_client.bucket.get_object(download_path)
            # 将内容写入临时文件
            with open(temp_path, 'wb') as f:
                f.write(result.read())
            logger.info(f"文件下载成功: {temp_path}")
        except Exception as e:
            raise Exception(f"从OSS下载文件失败: {download_path}, 错误: {str(e)}")
        
        # 3. 根据类型解析文件
        text_content = parse_file_by_type(temp_path, resource_type)
        
        # 4. 删除临时文件
        try:
            os.unlink(temp_path)
            logger.info(f"临时文件已删除: {temp_path}")
        except Exception as e:
            logger.warning(f"删除临时文件失败: {e}")
        
        # 5. 限制文本长度
        if len(text_content) > MAX_TEXT_LENGTH:
            logger.info(f"文本内容过长({len(text_content)}字符)，截取前{MAX_TEXT_LENGTH}字符")
            text_content = text_content[:MAX_TEXT_LENGTH] + "\n\n...(内容过长，已截取)"
        
        return text_content
    
    except Exception as e:
        logger.error(f"解析资源文件失败: {e}")
        raise Exception(f"解析资源文件失败: {str(e)}")


def get_file_suffix(resource_type: str, original_filename: str) -> str:
    """获取文件后缀"""
    type_mapping = {
        'pdf': '.pdf',
        'word': '.docx',
        'ppt': '.pptx',
        'excel': '.xlsx',
        'markdown': '.md',
        'text': '.txt'
    }
    
    suffix = type_mapping.get(resource_type)
    if not suffix and original_filename:
        # 从原始文件名提取后缀
        _, ext = os.path.splitext(original_filename)
        suffix = ext if ext else '.txt'
    
    return suffix or '.txt'


def parse_file_by_type(file_path: str, resource_type: str) -> str:
    """
    根据文件类型解析内容
    
    Args:
        file_path: 本地文件路径
        resource_type: 资源类型
        
    Returns:
        str: 解析后的文本内容
    """
    try:
        if resource_type == 'pdf':
            return parse_pdf(file_path)
        elif resource_type == 'word':
            return parse_word(file_path)
        elif resource_type == 'ppt':
            return parse_ppt(file_path)
        elif resource_type in ['markdown', 'text']:
            return parse_text(file_path)
        else:
            logger.warning(f"不支持的资源类型: {resource_type}，尝试作为文本解析")
            return parse_text(file_path)
    except Exception as e:
        logger.error(f"解析文件失败({resource_type}): {e}")
        return f"[文件解析失败: {str(e)}]"


def parse_pdf(file_path: str) -> str:
    """解析PDF文件"""
    try:
        from PyPDF2 import PdfReader
        
        reader = PdfReader(file_path)
        text_parts = []
        
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text.strip():
                text_parts.append(f"=== 第{i+1}页 ===\n{page_text}\n")
        
        return "\n".join(text_parts)
    
    except ImportError:
        logger.error("PyPDF2未安装，尝试使用pdfplumber")
        try:
            import pdfplumber
            
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"=== 第{i+1}页 ===\n{page_text}\n")
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"pdfplumber解析失败: {e}")
            return f"[PDF解析失败: 缺少必要的库]"
    
    except Exception as e:
        logger.error(f"PDF解析失败: {e}")
        return f"[PDF解析失败: {str(e)}]"


def parse_word(file_path: str) -> str:
    """解析Word文件"""
    try:
        from docx import Document
        
        doc = Document(file_path)
        text_parts = []
        
        # 解析段落
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # 解析表格
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        return "\n\n".join(text_parts)
    
    except ImportError:
        logger.error("python-docx未安装")
        return "[Word文件解析失败: 缺少必要的库python-docx]"
    except Exception as e:
        logger.error(f"Word解析失败: {e}")
        return f"[Word文件解析失败: {str(e)}]"


def parse_ppt(file_path: str) -> str:
    """解析PPT文件"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = []
        
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            
            # 提取幻灯片中的所有文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            
            if slide_texts:
                text_parts.append(f"=== 幻灯片{i+1} ===\n" + "\n".join(slide_texts))
        
        return "\n\n".join(text_parts)
    
    except ImportError:
        logger.error("python-pptx未安装")
        return "[PPT文件解析失败: 缺少必要的库python-pptx]"
    except Exception as e:
        logger.error(f"PPT解析失败: {e}")
        return f"[PPT文件解析失败: {str(e)}]"


def parse_text(file_path: str) -> str:
    """解析文本文件"""
    try:
        # 尝试多种编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                    return content
            except UnicodeDecodeError:
                continue
        
        # 所有编码都失败
        logger.error("无法识别文件编码")
        return "[文本文件编码无法识别]"
    
    except Exception as e:
        logger.error(f"文本解析失败: {e}")
        return f"[文本文件解析失败: {str(e)}]"
